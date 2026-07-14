"""PPTX report generator — 4-5 slide deck."""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from casmd.report.data import ReportData


def generate_pptx(data: ReportData, output_path: Path) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    pres = Presentation()
    blank_layout = pres.slide_layouts[6]  # blank

    # ---- Slide 1: Title ----
    s = pres.slides.add_slide(pres.slide_layouts[0])
    s.shapes.title.text = f"CasMD report — {data.job_name}"
    if len(s.placeholders) > 1:
        s.placeholders[1].text = f"Production: {data.production_ns:.0f} ns"

    # ---- Slide 2: Prediction summary (if present) ----
    if data.prediction is not None:
        s = pres.slides.add_slide(pres.slide_layouts[1])
        s.shapes.title.text = "Structure prediction"
        body = s.placeholders[1].text_frame
        body.text = f"Backend: {data.prediction.backend} (model {data.prediction.model_id})"
        for label, val in [
            ("iPTM", data.prediction.iptm),
            ("pTM", data.prediction.ptm),
            ("pLDDT (mean)", data.prediction.plddt_mean),
        ]:
            p = body.add_paragraph()
            p.text = f"{label}: {val:.3f}" if val is not None else f"{label}: --"

    # ---- Slide 3: Simulation summary ----
    s = pres.slides.add_slide(pres.slide_layouts[1])
    s.shapes.title.text = "Simulation summary"
    body = s.placeholders[1].text_frame
    body.text = f"Frames: {data.analysis.n_frames}; equilibration discard: {data.analysis.equil_skip} frames"
    for label, val in [
        ("Protein Cα RMSD (equil. mean)", data.analysis.protein_rmsd_equil_mean_A),
        ("Protein Cα RMSF (mean)", data.analysis.protein_rmsf_mean_A),
        ("Protein Rg (equil. mean)", data.analysis.protein_rg_equil_mean_A),
    ]:
        p = body.add_paragraph()
        p.text = f"{label}: {val:.2f} Å" if val is not None else f"{label}: --"

    # ---- Slide 4: RMSD + RMSF figures ----
    s = pres.slides.add_slide(blank_layout)
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title.text_frame.text = "Protein dynamics"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    if "rmsd" in data.figures and data.figures["rmsd"].exists():
        s.shapes.add_picture(str(data.figures["rmsd"]),
                              Inches(0.3), Inches(1.0), width=Inches(4.7))
    if "rmsf" in data.figures and data.figures["rmsf"].exists():
        s.shapes.add_picture(str(data.figures["rmsf"]),
                              Inches(5.0), Inches(1.0), width=Inches(4.7))

    # ---- Slide 5: Rg + Interpretation ----
    s = pres.slides.add_slide(blank_layout)
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title.text_frame.text = "Compactness & interpretation"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    if "rg" in data.figures and data.figures["rg"].exists():
        s.shapes.add_picture(str(data.figures["rg"]),
                              Inches(0.3), Inches(1.0), width=Inches(4.7))
    interp = s.shapes.add_textbox(Inches(5.0), Inches(1.0), Inches(4.7), Inches(5.0))
    tf = interp.text_frame
    tf.word_wrap = True
    tf.text = data.interpretation

    pres.save(str(output_path))
    return output_path


def generate_comparison_pptx(cd, out_path) -> None:
    """Write a comparison PPTX deck."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pathlib import Path as _Path

    prs = Presentation()

    # Slide 1 — title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = f"{cd.label_a} vs {cd.label_b}"
    if s.placeholders[1].has_text_frame:
        s.placeholders[1].text = "CasMD comparison report · v0.9.1"

    # Slide 2 — summary table
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Summary"
    rows = len(cd.stats) + 1
    table = s.shapes.add_table(rows=rows, cols=4,
                                left=Inches(0.5), top=Inches(1.5),
                                width=Inches(9), height=Inches(3.5)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = cd.label_a
    table.cell(0, 2).text = cd.label_b
    table.cell(0, 3).text = "Δ"
    for i, (key, stat) in enumerate(cd.stats.items(), start=1):
        table.cell(i, 0).text = key.upper()
        table.cell(i, 1).text = f"{stat.mean_a:.2f}"
        table.cell(i, 2).text = f"{stat.mean_b:.2f}"
        table.cell(i, 3).text = f"{stat.delta:+.2f}"

    # Slide 3 — figures (if any)
    if cd.figures_a or cd.figures_b:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "Dynamics"
        for i, key in enumerate(("rmsd", "rmsf", "rg")):
            if key in cd.figures_a:
                s.shapes.add_picture(str(cd.figures_a[key]),
                                      left=Inches(0.2 + i * 3.3),
                                      top=Inches(1.5),
                                      width=Inches(3.0))
            if key in cd.figures_b:
                s.shapes.add_picture(str(cd.figures_b[key]),
                                      left=Inches(0.2 + i * 3.3),
                                      top=Inches(4.0),
                                      width=Inches(3.0))

    # Slide N — interpretation
    if cd.interpretation:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "Interpretation"
        # Use a textbox (more reliable than searching for body placeholders).
        box = s.shapes.add_textbox(Inches(0.5), Inches(1.5),
                                    Inches(9), Inches(5))
        box.text_frame.text = cd.interpretation

    prs.save(str(_Path(out_path)))
